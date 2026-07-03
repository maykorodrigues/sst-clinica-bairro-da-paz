"""
Gerador RMAR SST Card — Bairro da Paz
Adapta o template original (Cidade Ademar) para a realidade SST Card Bairro da Paz.
Uso: python gerar_rmar_sst.py
Saída: RMAR-SST-Card-Bairro-da-Paz-[MES].pptx
"""

import os, sys, copy, shutil
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import ChartData
from lxml import etree

# ── Configuração do relatório ────────────────────────────────────────────────
MES_REFERENCIA = "06/2026"
UNIDADE        = "Bairro da Paz, Salvador/BA"
PREPARADO_POR  = "Mayko Rodrigues | Consultor Comercial"
NOME_SAIDA     = "RMAR-SST-Card-Bairro-da-Paz-ABRIL-MAIO-JUNHO-2026.pptx"

TEMPLATE = os.path.join(os.path.dirname(__file__), "rmar-temp.pptx")
SAIDA    = os.path.join(os.path.dirname(__file__), NOME_SAIDA)

# ── Dados históricos — tabelas de adimplência e churn (preencher mensalmente) ──
# QCA = Quantidade de Clientes Ativos | QIA = Quantidade de Itens Arrecadados (pagantes)
# Adimplência: cada linha = [MÊS, QCA, QIA, ADIMPLÊNCIA %]
#
# Abril/2026: base total SST Card pré-advisory (fonte: Boom, Karine 17/04/2026)
#   376 titulares | 41 inadimplentes → 335 pagantes | adimplência 89,10%
# Maio/2026: cohort advisory (novas adesões + Tenex reativados via campanha Perdão de Dívida)
#   35 membros do período advisory | 28 pagantes | adimplência 80,00%
ADIMPLENCIA = [
    ["04/2026", "376", "335", "89,10%"],
    ["05/2026",  "35",  "28", "80,00%"],
    # Junho/2026: base RECORRENTE auditada no Asaas (extrato 01-30/06)
    #   327 contratos recorrentes contratados | 178 pagantes no mês → 54,44%
    #   (R$ 7,9k/mês não arrecadados = alvo da Máquina de Cobrança)
    ["06/2026", "327", "178", "54,44%"],
    # formato: ["MM/AAAA", "total_ativos", "pagantes_no_mes", "percentual"]
]

# Churn: cada linha = [MÊS, QCA, DESFILIAÇÕES, CHURN RATE]
# Abril: 4 cancelamentos (Boom 17/04) | Churn rate 1,06%
# Maio: 0 cancelamentos no cohort advisory
CHURN = [
    ["04/2026", "376",  "4", "1,06%"],
    ["05/2026",  "35",  "0", "0,00%"],
    ["06/2026", "327",  "6", "1,83%"],  # junho — estimativa (confirmar cancelamentos com Karine)
]

# ── Dados dos gráficos — Abril e Maio/2026 ──────────────────────────────────
# Slide 5: Adesões Totais (Novas + Reativações)
# Slide 6: Novas Adesões
# Slide 7: Reativações Tenex
# Slide 8: Canal de Captação (stacked: WhatsApp+Insta vs Presencial)
# Slide 9: Composição por Canal (100% stacked: 4 canais)
# Slide 11: Evolução Receita Mensal (R$)
# Slide 12: Composição da Receita por Plano (stacked 4 categorias)
# Slide 13: Taxa de Adesão Arrecadada (R$)
#
# Metodologia abril/2026 (linha de base pré-advisory):
#  Fonte: relato Karine no grupo WhatsApp 17/04/2026 (sistema Boom)
#  - 24 novas adesões orgânicas (sem campanha Tenex)
#  - 0 reativações Tenex (campanha Perdão de Dívida ainda não iniciada)
#  - Canal: Presencial/Klingo 20, WhatsApp 3, Instagram 1
#  - Receita (estimativa mix):
#      18 × R$39,90  = R$718  (Individual Plano Prata)
#       6 × R$65,00  = R$390  (Família)
#      20 × R$25,00  = R$500  (Taxa de Adesão — pré-nova precificação 02/06)
#       0 × R$0      = R$0    (sem Tenex)
#      TOTAL ≈ R$1.608
#
# Metodologia maio/2026 (Mês 1 do advisory):
#  - 25 novas adesões confirmadas (mix Individual + Família, maio todo)
#  - 10 reativações Tenex pagas (campanha Perdão de Dívida)
#  - Canal: WhatsApp/Instagram 17, Presencial 8 (das 25 novas)
#  - Raquel (marketing): 370 pessoas faladas, 50 interessados, 14,7k views Instagram
#  - Receita:  15 × R$34,90 = R$524  (Individual, pré-22/05)
#              8 × R$65     = R$520  (Família, pré-22/05)
#              2 × R$39,90  = R$80   (Individual, pós-22/05 — incluso no total)
#              5 × R$40     = R$200  (taxa de adesão, nova precificação pós-22/05)
#             10 × R$34,90  = R$349  (reativações Tenex)
#             TOTAL ≈ R$1.673
# Junho/2026 (Mês 2 do advisory) — dados AUDITADOS no extrato Asaas (01-30/06):
#  - Receita total (produção): R$ 12.144 = Asaas R$ 9.811 (recorrente cartão R$ 3.794
#    + ativo Pix R$ 5.927 + débito R$ 90) + Dinheiro recepção R$ 2.332 (BOON+SST Card)
#  - Aporte separado: R$ 15.000 (Maria Adélia) — NÃO entra na receita
#  ⚠️ Adesões/reativações/canais de junho = ESTIMATIVAS — confirmar com Karine
GRAFICOS = {
    # slide_idx (0-based): {'categories': [...], 'series': {nome: [v_abr, v_mai, v_jun]}}
    4:  {  # Slide 5 — Adesões Totais
        'categories': ['04/2026', '05/2026', '06/2026'],
        'series': {'Adesões Totais (Novas + Reativações)': [24, 35, 42]},
    },
    5:  {  # Slide 6 — Novas Adesões
        'categories': ['04/2026', '05/2026', '06/2026'],
        'series': {'Novas Adesões SST Card': [24, 25, 30]},
    },
    6:  {  # Slide 7 — Reativações Tenex
        'categories': ['04/2026', '05/2026', '06/2026'],
        'series': {'Reativações Tenex': [0, 10, 12]},
    },
    7:  {  # Slide 8 — Canal de Captação (stacked)
        'categories': ['04/2026', '05/2026', '06/2026'],
        'series': {'WhatsApp / Instagram': [4, 17, 22], 'Presencial (Campo)': [20, 8, 20]},
    },
    8:  {  # Slide 9 — Composição por Canal (100% stacked)
        'categories': ['04/2026', '05/2026', '06/2026'],
        'series': {'WhatsApp': [3, 15, 18], 'Presencial': [20, 8, 20], 'Instagram': [1, 2, 4], 'Indicação': [0, 0, 0]},
    },
    10: {  # Slide 11 — Receita Total / Produção (R$) — AUDITADA
        'categories': ['04/2026', '05/2026', '06/2026'],
        'series': {'Receita Total / Produção (R$)': [5929, 10127, 12144]},
    },
    11: {  # Slide 12 — Composição da Receita por Fonte (stacked)
        'categories': ['04/2026', '05/2026', '06/2026'],
        'series': {
            'Recorrente (cartão)':   [2000, 3000, 3794],
            'Ativo (Pix/cobrança)':  [3929, 4222, 5927],
            'Dinheiro (recepção)':   [   0, 2905, 2332],
            'Débito/Outros':         [   0,    0,   90],
        },
    },
    12: {  # Slide 13 — Taxa de Adesão Arrecadada (R$)
        'categories': ['04/2026', '05/2026', '06/2026'],
        'series': {'Taxa de Adesão (R$)': [500, 200, 1050]},
    },
}

# ── Renomeação de séries nos gráficos (por índice de slide, 0-based) ─────────
SERIES_RENAMES = {
    4:  {0: "Novas Adesões + Reativações"},                          # slide 5
    5:  {0: "Novas Adesões"},                                        # slide 6
    6:  {0: "Reativações Tenex"},                                    # slide 7
    7:  {0: "WhatsApp / Instagram", 1: "Presencial (Campo)"},        # slide 8
    8:  {0: "WhatsApp", 1: "Presencial", 2: "Instagram", 3: "Indicação"},  # slide 9
    10: {0: "Receita Total / Produção (R$)"},                        # slide 11
    11: {0: "Recorrente (cartão)", 1: "Ativo (Pix/cobrança)", 2: "Dinheiro (recepção)", 3: "Débito/Outros"},  # slide 12
    12: {0: "Taxa de Adesão (R$)"},                                  # slide 13
}

# ── Substituições de texto por slide (0-based) ───────────────────────────────
TEXT_REPLACEMENTS = {
    0: [                                              # slide 1
        ("10/2020", MES_REFERENCIA),
        ("10/2021", MES_REFERENCIA),
    ],
    1: [                                              # slide 2
        # específicos primeiro (mais longos antes de substrings)
        ("Equipe Cidade Ademar", PREPARADO_POR),
        ("cidade ademar/sp", UNIDADE),
        ("CIDADE ADEMAR/SP", UNIDADE),
        ("Cidade Ademar/SP", UNIDADE),
        ("Cidade Ademar", UNIDADE),
    ],
    3: [                                              # slide 4
        ("10/2020", MES_REFERENCIA),
    ],
    4: [                                              # slide 5
        # texto tem quebra de linha \x0b no original — substituir e limpar subtitle
        ("Prospecções total", "Adesões Totais (Novas + Reativações Tenex)"),
        # limpar o texto que fica após o \x0b (quebra de linha interna)
        ("\x0b(novas prospecções + refiliações)", ""),
        ("(novas prospecções + refiliações)", ""),
    ],
    5: [                                              # slide 6
        ("Novas prospecções", "Novas Adesões SST Card"),
    ],
    6: [                                              # slide 7
        ("Refiliações", "Reativações Tenex"),
    ],
    7: [                                              # slide 8
        ("Modalidade prospecções", "Canal de Captação"),
    ],
    8: [                                              # slide 9
        ("Modalidade prospecções e refiliações", "Composição por Canal de Venda"),
    ],
    10: [                                             # slide 11
        ("Evolução qia – QUANTIDADE DE ITENS ARRECADADOS",
         "Evolução Receita Mensal — SST Card (R$)"),
        ("Evolução qia", "Evolução Receita (R$)"),
    ],
    11: [                                             # slide 12
        ("Composição qia", "Composição da Receita por Plano"),
        ("Composição QIA", "Composição da Receita por Plano"),
    ],
    12: [                                             # slide 13
        ("HOMOLOGAÇÃO", "TAXA DE ADESÃO ARRECADADA (R$)"),
        ("Homologação", "Taxa de Adesão (R$)"),
    ],
}

# ── Helpers ──────────────────────────────────────────────────────────────────

def replace_para(para, old, new):
    """Substitui texto em um parágrafo preservando formatação do primeiro run.
    Trunca em \x0b (quebra de linha XML) para evitar _x000B_ no output."""
    full = para.text
    if old.lower() not in full.lower():
        return False
    idx = full.lower().find(old.lower())
    new_full = full[:idx] + new + full[idx + len(old):]
    # Truncar na primeira quebra de linha interna para não poluir com _x000B_
    new_first_line = new_full.split('\x0b')[0].strip()
    if para.runs:
        para.runs[0].text = new_first_line
        for run in para.runs[1:]:
            run.text = ""
    return True


def replace_shape_text(shape, replacements):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for old, new in replacements:
                replace_para(para, old, new)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for old, new in replacements:
                        replace_para(para, old, new)


def remove_line_breaks(slide):
    """Remove elementos <a:br/> dos text frames (evita _x000B_ no output)."""
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for br_el in para._p.findall(f"{{{NS_A}}}br"):
                para._p.remove(br_el)


def clear_all_tables_in_slide(slide, skip_first_table_header=True):
    """Limpa TODAS as tabelas de um slide (inclusive tabelas 1x4 separadas).
    skip_first_table_header: preserva header da tabela principal; tabelas menores
    (1 linha) são limpas por completo."""
    first_table = True
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        # Pula cabeçalho só na tabela principal (primeira, com >1 linha)
        skip = skip_first_table_header and first_table and len(tbl.rows) > 1
        start = 1 if skip else 0
        first_table = False
        for r_idx in range(start, len(tbl.rows)):
            for c_idx in range(len(tbl.columns)):
                set_cell_text(tbl.cell(r_idx, c_idx), "—")


def rename_chart_series(slide, series_map):
    """Renomeia séries de gráficos via XML (suporte a cache de string)."""
    NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    for shape in slide.shapes:
        try:
            chart = shape.chart
        except Exception:
            continue
        for idx, new_name in series_map.items():
            try:
                ser = chart.series[idx]
                # tenta via cache de string no XML
                ser_el = ser._element
                # <c:tx><c:strRef><c:strCache><c:pt idx="0"><c:v>Nome</c:v>
                for v_el in ser_el.iter(f"{{{NS}}}v"):
                    if v_el.getparent().getparent().tag == f"{{{NS}}}pt":
                        gp = v_el.getparent().getparent().getparent()
                        if gp.tag == f"{{{NS}}}strCache":
                            v_el.text = new_name
                            print(f"  série {idx} → '{new_name}'")
                            break
                # fallback: fórmula direta
                for f_el in ser_el.iter(f"{{{NS}}}f"):
                    if "Series" in (f_el.text or ""):
                        f_el.text = f'"{new_name}"'
            except Exception as e:
                print(f"  AVISO: série {idx} — {e}")


def update_chart_data(slide, categories, series_data):
    """Substitui os dados do primeiro gráfico do slide usando ChartData."""
    for shape in slide.shapes:
        try:
            chart = shape.chart
        except Exception:
            continue
        cd = ChartData()
        cd.categories = categories
        for name, values in series_data.items():
            cd.add_series(name, tuple(values))
        chart.replace_data(cd)
        names = list(series_data.keys())
        sys.stdout.buffer.write(f"  Grafico atualizado: {names}\n".encode('utf-8'))
        return


def set_cell_text(cell, val):
    """Define texto de uma célula de tabela preservando formatação."""
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = val
        for run in para.runs[1:]:
            run.text = ""
    else:
        from pptx.util import Pt
        run = para.add_run()
        run.text = val


def fill_table(slide, data, start_row=1):
    """Preenche a primeira tabela do slide: insere dados novos e limpa linhas antigas."""
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            n_rows = len(tbl.rows)
            # Preencher com dados novos
            for r_idx, row_data in enumerate(data):
                r = start_row + r_idx
                if r >= n_rows:
                    break
                for c_idx, val in enumerate(row_data):
                    if c_idx < len(tbl.columns):
                        set_cell_text(tbl.cell(r, c_idx), val)
            # Limpar linhas restantes (dados antigos da Cidade Ademar)
            clear_from = start_row + len(data)
            for r in range(clear_from, n_rows):
                for c_idx in range(len(tbl.columns)):
                    set_cell_text(tbl.cell(r, c_idx), "—")
            return  # apenas primeira tabela


# ── Pipeline principal ────────────────────────────────────────────────────────

def main():
    if not os.path.exists(TEMPLATE):
        sys.exit(f"ERRO: template não encontrado em {TEMPLATE}")

    prs = Presentation(TEMPLATE)
    slides = list(prs.slides)

    print("Processando slides...")

    # 1. Substituições de texto
    for slide_idx, replacements in TEXT_REPLACEMENTS.items():
        if slide_idx >= len(slides):
            continue
        slide = slides[slide_idx]
        for shape in slide.shapes:
            replace_shape_text(shape, replacements)
        print(f"  Slide {slide_idx+1}: texto atualizado")

    # 2. Renomear séries nos gráficos
    print("Renomeando séries...")
    for slide_idx, series_map in SERIES_RENAMES.items():
        if slide_idx >= len(slides):
            continue
        rename_chart_series(slides[slide_idx], series_map)

    # 3. Atualizar dados dos gráficos com valores reais de maio/2026
    sys.stdout.buffer.write(b"Atualizando graficos...\n")
    for slide_idx, dados in GRAFICOS.items():
        if slide_idx >= len(slides):
            continue
        sys.stdout.buffer.write(f"  Slide {slide_idx+1}:\n".encode('utf-8'))
        update_chart_data(slides[slide_idx], dados['categories'], dados['series'])
    sys.stdout.buffer.flush()

    # 4. Remover quebras de linha internas (slide 5 — _x000B_)
    if len(slides) > 4:
        remove_line_breaks(slides[4])
        print("  Slide 5: quebras de linha removidas")

    # 5. Preencher tabelas — limpar TODOS os dados Cidade Ademar primeiro
    print("Preenchendo tabelas...")

    # Slides 14-15 (adimplência, índices 13-14): limpar todas as tabelas
    for idx in [13, 14]:
        if idx < len(slides):
            clear_all_tables_in_slide(slides[idx], skip_first_table_header=True)

    # Slides 17-18 (churn, índices 16-17): idem
    for idx in [16, 17]:
        if idx < len(slides):
            clear_all_tables_in_slide(slides[idx], skip_first_table_header=True)

    # Agora preencher com dados reais
    if len(slides) > 13 and ADIMPLENCIA:
        fill_table(slides[13], ADIMPLENCIA, start_row=1)
        print("  Slide 14: adimplência atualizada")

    if len(slides) > 16 and CHURN:
        fill_table(slides[16], CHURN, start_row=1)
        print("  Slide 17: churn rate atualizado")

    # 4. Salvar
    prs.save(SAIDA)
    sys.stdout.buffer.write(f"\nOK — Arquivo salvo em:\n   {SAIDA}\n".encode('utf-8'))
    sys.stdout.buffer.write(b"\nProximos passos:\n")
    sys.stdout.buffer.write(b"  1. Abrir no PowerPoint e atualizar graficos com dados reais\n")
    sys.stdout.buffer.write(b"  2. Slide 3: editar pilares do modelo (visual, nao texto)\n")
    sys.stdout.buffer.write(b"  3. Slides 5-9: adicionar dados mensais nos graficos\n")
    sys.stdout.buffer.write(b"  4. Slides 14-15: adicionar meses em ADIMPLENCIA no script\n")
    sys.stdout.buffer.write(b"  5. Slides 17-18: idem para CHURN no script\n")
    sys.stdout.buffer.flush()


if __name__ == "__main__":
    main()
